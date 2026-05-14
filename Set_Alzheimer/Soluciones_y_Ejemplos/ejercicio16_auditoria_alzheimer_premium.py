import os
import sys

import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.append(ROOT)
from alzheimer_utils import load_and_clean_data


WIDTH = Inches(10)
HEIGHT = Inches(7.5)
COLOR_BLUE = RGBColor(15, 32, 60)
COLOR_TEXT = RGBColor(40, 40, 40)


def style_text(frame, size=18, bold=False, color=COLOR_TEXT):
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for paragraph in frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.name = "Segoe UI"
        paragraph.font.color.rgb = color


def add_slide(prs, title, text, img_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    rect = slide.shapes.add_shape(6, 0, 0, Inches(0.15), HEIGHT)
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_BLUE
    rect.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.9))
    title_box.text = title
    style_text(title_box.text_frame, size=26, bold=True, color=COLOR_BLUE)

    if img_path and os.path.exists(img_path):
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(4.2), Inches(5.0))
        body.text = text
        style_text(body.text_frame, size=14)
        slide.shapes.add_picture(img_path, Inches(5.1), Inches(1.7), width=Inches(4.1))
    else:
        body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(5.0))
        body.text = text
        style_text(body.text_frame, size=16)


df = load_and_clean_data(os.path.join(ROOT, "base_datos_alzheimer.csv"))
sns.set_theme(style="whitegrid")

plt.figure(figsize=(6, 4))
sns.regplot(x="PET_Amiloide_Centiloid", y="MMSE_Actual", data=df, color="#1B4F72")
plt.tight_layout()
plot_a = os.path.join(ROOT, "plot_alz_a.png")
plt.savefig(plot_a, dpi=300, transparent=True)
plt.close()

plt.figure(figsize=(6, 4))
sns.boxplot(x="Tratamiento_Desc", y="Declive_MMSE_Anual", data=df, hue="Tratamiento_Desc", palette="Greens", legend=False)
plt.xticks(rotation=18)
plt.tight_layout()
plot_b = os.path.join(ROOT, "plot_alz_b.png")
plt.savefig(plot_b, dpi=300, transparent=True)
plt.close()

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "AUDITORIA CLINICA ALZHEIMER\nSELECCION DE CANDIDATOS ANTIAMILOIDE"
slide.shapes.placeholders[1].text = "Biomarcadores multimodales, riesgo radiologico y priorizacion clinica"

add_slide(prs, "Resumen ejecutivo", "- La elegibilidad antiamiloide depende de estadio, amiloide y seguridad MRI.\n- El set separa claramente comparadores biomarcador-negativos.\n- La RM de microhemorragias condiciona seleccion terapéutica.")
add_slide(prs, "Carga amiloide y cognicion", "La relacion entre centiloides y rendimiento cognitivo es monotona pero incompleta: la neurodegeneracion explica parte del deterioro clinico.", plot_a)
add_slide(prs, "Tratamiento y velocidad de declive", "La comparacion cruda por tratamiento esta influida por sesgo de indicacion. Los pacientes con memantina o combinacion suelen estar mas avanzados.", plot_b)

for title in [
    "ATN y seleccion terapeutica",
    "Riesgo de ARIA y microhemorragias",
    "Biomarcadores plasmaticos como triaje",
    "Concordancia LCR-PET-RM",
    "ApoE4 y estratificacion de riesgo",
    "Lectura funcional con FAQ y CDR-SB",
    "DCL amnesico frente a demencia leve",
    "Comparadores biomarcador-negativos",
    "Implicaciones para circuitos asistenciales",
    "Priorizacion de recursos PET/LCR",
]:
    add_slide(prs, title, "Hallazgos sinteticos de la cohorte para discusion docente y construccion de circuitos clinicos reproducibles.")

add_slide(prs, "Recomendaciones", "1. Confirmar amiloide antes de discutir anticuerpos.\n2. Integrar MRI de seguridad antes de decidir tratamiento.\n3. No interpretar comparaciones observacionales sin ajustar por estadio.")
add_slide(prs, "Cierre", "La cohorte sintetica permite entrenar razonamiento clinico, analitica y comunicacion ejecutiva sobre Alzheimer temprano.")

output = os.path.join(ROOT, "Auditoria_Clinica_Alzheimer.pptx")
prs.save(output)
print(f"Presentacion generada: {output}")
