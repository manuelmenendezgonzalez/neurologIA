import os
import sys

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.append(ROOT)
from alzheimer_utils import load_and_clean_data


COLOR_NAVY = RGBColor(15, 32, 60)
COLOR_TEAL = RGBColor(25, 111, 61)
COLOR_GOLD = RGBColor(196, 155, 45)
COLOR_WHITE = RGBColor(255, 255, 255)
STAGE_ORDER = [
    "Queja subjetiva / No demencia",
    "DCL amnésico",
    "Demencia leve",
    "Demencia moderada",
]


def apply_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_NAVY


def add_title(slide, title):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.7))
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = COLOR_GOLD


def add_text_image_slide(prs, title, bullets, img_path=None, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    add_title(slide, title)

    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.25), Inches(5.7), Inches(5.4))
    tf = body.text_frame
    tf.clear()
    first = True
    for bullet in bullets:
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = COLOR_WHITE
        first = False

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(6.55), Inches(1.35), width=Inches(6.0))

    if footer:
        footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.4))
        p = footer_box.text_frame.paragraphs[0]
        p.text = footer
        p.font.size = Pt(10)
        p.font.color.rgb = COLOR_WHITE


df = load_and_clean_data(os.path.join(ROOT, "base_datos_alzheimer.csv"))
df["Estado_Progresion"] = np.where(
    (df["Declive_MMSE_Anual"] >= 2.0) | (df["Cambio_CDRSB_Anual"] >= 0.9),
    "Progresor rapido",
    "Lento/estable",
)

stage_summary = (
    df.groupby("Estadio_Desc")
    .agg(
        n=("ID", "count"),
        mmse=("MMSE_Actual", "median"),
        declive=("Declive_MMSE_Anual", "median"),
        cdr=("CDR_SB_Actual", "median"),
        faq=("FAQ_Actual", "median"),
    )
    .reindex(STAGE_ORDER)
    .round(2)
)
fast_rate = (
    df.groupby("Estadio_Desc")["Estado_Progresion"]
    .apply(lambda s: (s == "Progresor rapido").mean() * 100)
    .reindex(STAGE_ORDER)
    .round(1)
)
stage_summary["fast_rate"] = fast_rate

total = len(df)
fast = int((df["Estado_Progresion"] == "Progresor rapido").sum())
amyloid_positive = int(df["Amyloid_Pos"].sum())
antiamyloid_candidates = int(df["Candidato_Antiamiloide"].sum())
antiamyloid_treated = int(df["Tratamiento_Cod"].isin([6, 7]).sum())
early_ad = df[(df["Estadio_Cod"].isin([1, 2])) & (df["Amyloid_Pos"] == 1)]
treated_early_ad = early_ad[early_ad["Tratamiento_Cod"].isin([6, 7])]
untreated_early_ad = early_ad[~early_ad["Tratamiento_Cod"].isin([6, 7])]

atn_top = (
    df["ATN_Perfil"]
    .value_counts()
    .head(3)
    .reset_index()
    .apply(lambda row: f"{row['ATN_Perfil']}: {row['count']} pacientes", axis=1)
    .tolist()
)

treatment_counts = df["Tratamiento_Desc"].value_counts()
main_treatments = treatment_counts[treatment_counts >= 5].index.tolist()
treatment_summary = (
    df[df["Tratamiento_Desc"].isin(main_treatments)]
    .groupby("Tratamiento_Desc")
    .agg(
        n=("ID", "count"),
        declive=("Declive_MMSE_Anual", "mean"),
        cdr=("Cambio_CDRSB_Anual", "mean"),
    )
    .sort_values("declive")
    .round(2)
)

candidate_stage = (
    df.groupby("Estadio_Desc")
    .agg(
        candidatos=("Candidato_Antiamiloide", "sum"),
        tratados=("Tratamiento_Cod", lambda s: int(s.isin([6, 7]).sum())),
    )
    .reindex(STAGE_ORDER)
    .reset_index()
)
candidate_long = candidate_stage.melt(
    id_vars="Estadio_Desc",
    value_vars=["candidatos", "tratados"],
    var_name="Grupo",
    value_name="Pacientes",
)
stage_distribution_text = ", ".join(
    [f"{stage}={int(stage_summary.loc[stage, 'n'])}" for stage in STAGE_ORDER]
)

cor_amyloid_declive = df["PET_Amiloide_Centiloid"].corr(df["Declive_MMSE_Anual"])
cor_fdg_cdr = df["FDG_PET_TP_SUVR"].corr(df["CDR_SB_Actual"])
cor_hippo_cdr = df["RM_Hipocampo_Pctl"].corr(df["CDR_SB_Actual"])

sns.set_theme(style="whitegrid")

plot_stage = os.path.join(ROOT, "ej2_alz_stage_declive.png")
plt.figure(figsize=(8.0, 5.6))
ax = sns.boxplot(
    data=df,
    x="Estadio_Desc",
    y="Declive_MMSE_Anual",
    order=STAGE_ORDER,
    hue="Estado_Progresion",
    palette=["#5DADE2", "#F5B041"],
)
if ax.legend_:
    ax.legend_.set_title("")
plt.xlabel("")
plt.ylabel("Declive anual MMSE")
plt.xticks(rotation=12)
plt.tight_layout()
plt.savefig(plot_stage, dpi=300)
plt.close()

plot_function = os.path.join(ROOT, "ej2_alz_functional_change.png")
plt.figure(figsize=(8.0, 5.6))
sns.boxplot(
    data=df,
    x="Estadio_Desc",
    y="Cambio_CDRSB_Anual",
    order=STAGE_ORDER,
    color="#52BE80",
)
plt.xlabel("")
plt.ylabel("Cambio anual CDR-SB")
plt.xticks(rotation=12)
plt.tight_layout()
plt.savefig(plot_function, dpi=300)
plt.close()

plot_biomarkers = os.path.join(ROOT, "ej2_alz_biomarkers.png")
plt.figure(figsize=(8.0, 5.6))
sns.scatterplot(
    data=df,
    x="PET_Amiloide_Centiloid",
    y="Declive_MMSE_Anual",
    hue="Estadio_Desc",
    hue_order=STAGE_ORDER,
    palette="viridis",
    s=70,
)
plt.xlabel("PET amiloide (centiloid)")
plt.ylabel("Declive anual MMSE")
plt.tight_layout()
plt.savefig(plot_biomarkers, dpi=300)
plt.close()

plot_treatment = os.path.join(ROOT, "ej2_alz_treatment.png")
plt.figure(figsize=(8.4, 5.6))
sns.boxplot(
    data=df[df["Tratamiento_Desc"].isin(main_treatments)],
    x="Tratamiento_Desc",
    y="Declive_MMSE_Anual",
    order=treatment_summary.index.tolist(),
    color="#AF7AC5",
)
plt.xlabel("")
plt.ylabel("Declive anual MMSE")
plt.xticks(rotation=22)
plt.tight_layout()
plt.savefig(plot_treatment, dpi=300)
plt.close()

plot_candidates = os.path.join(ROOT, "ej2_alz_candidates.png")
plt.figure(figsize=(8.0, 5.6))
sns.barplot(
    data=candidate_long,
    x="Estadio_Desc",
    y="Pacientes",
    hue="Grupo",
    order=STAGE_ORDER,
    palette=["#52BE80", "#F4D03F"],
)
plt.xlabel("")
plt.ylabel("Pacientes")
plt.xticks(rotation=12)
plt.tight_layout()
plt.savefig(plot_candidates, dpi=300)
plt.close()

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

cover = prs.slides.add_slide(prs.slide_layouts[6])
apply_background(cover)
title_box = cover.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.8), Inches(1.2))
title_p = title_box.text_frame.paragraphs[0]
title_p.text = "Ejercicio 2: Progresion Clinica en Alzheimer"
title_p.alignment = PP_ALIGN.CENTER
title_p.font.size = Pt(34)
title_p.font.bold = True
title_p.font.color.rgb = COLOR_GOLD

subtitle_box = cover.shapes.add_textbox(Inches(1.4), Inches(3.0), Inches(10.5), Inches(1.0))
subtitle_p = subtitle_box.text_frame.paragraphs[0]
subtitle_p.text = "Cohorte, velocidad de deterioro, biomarcadores y lectura terapeutica"
subtitle_p.alignment = PP_ALIGN.CENTER
subtitle_p.font.size = Pt(20)
subtitle_p.font.color.rgb = COLOR_WHITE

add_text_image_slide(
    prs,
    "Resumen clinico",
    [
        f"Cohorte total: {total} pacientes; amiloide positivo en {amyloid_positive} ({amyloid_positive / total * 100:.1f}%).",
        f"Progresores rapidos: {fast} ({fast / total * 100:.1f}%) usando MMSE anual >= 2.0 o CDR-SB anual >= 0.9.",
        f"Candidatos antiamiloide: {antiamyloid_candidates}; tratados con antiamiloide: {antiamyloid_treated}.",
        "La cohorte separa fases preclinicas, DCL amnesico, demencia leve y demencia moderada con biomarcadores concordantes.",
    ],
    footer="Esta diapositiva debe convertirse en el resumen ejecutivo del manuscrito.",
)

add_text_image_slide(
    prs,
    "Cohorte y estratificacion",
    [
        f"Distribucion por estadio: {stage_distribution_text}.",
        f"ATN predominante: {atn_top[0]}; secundarios: {atn_top[1]} y {atn_top[2]}.",
        "La mediana de MMSE desciende de forma monotona entre estadio 0, DCL, demencia leve y moderada.",
        "La discapacidad funcional aumenta en paralelo y evita la mezcla clinica que tenia la cohorte previa.",
    ],
    img_path=plot_stage,
    footer="En el articulo, esta diapositiva debe ampliarse como seccion de descripcion de cohorte y estratificacion.",
)

add_text_image_slide(
    prs,
    "Progresion cognitiva anual",
    [
        f"Declive anual mediano de MMSE: DCL {stage_summary.loc['DCL amnésico', 'declive']}, demencia leve {stage_summary.loc['Demencia leve', 'declive']}, demencia moderada {stage_summary.loc['Demencia moderada', 'declive']}.",
        f"Progresores rapidos por estadio: DCL {stage_summary.loc['DCL amnésico', 'fast_rate']}%, demencia leve {stage_summary.loc['Demencia leve', 'fast_rate']}%, demencia moderada {stage_summary.loc['Demencia moderada', 'fast_rate']}%.",
        "El estadio 0 conserva declive bajo y MMSE alto, coherente con queja subjetiva o fase preclinica.",
        "La velocidad de deterioro debe interpretarse junto al estadio basal y no como un marcador aislado.",
    ],
    img_path=plot_stage,
    footer="Cada uno de estos hallazgos debe desarrollarse en un apartado propio del manuscrito.",
)

add_text_image_slide(
    prs,
    "Progresion funcional anual",
    [
        f"CDR-SB anual mediano: DCL {df.loc[df['Estadio_Desc']=='DCL amnésico', 'Cambio_CDRSB_Anual'].median():.2f}, demencia leve {df.loc[df['Estadio_Desc']=='Demencia leve', 'Cambio_CDRSB_Anual'].median():.2f}, demencia moderada {df.loc[df['Estadio_Desc']=='Demencia moderada', 'Cambio_CDRSB_Anual'].median():.2f}.",
        f"FAQ mediano por estadio: 0={stage_summary.loc['Queja subjetiva / No demencia', 'faq']}, 1={stage_summary.loc['DCL amnésico', 'faq']}, 2={stage_summary.loc['Demencia leve', 'faq']}, 3={stage_summary.loc['Demencia moderada', 'faq']}.",
        "La funcionalidad acompana a la carga neurodegenerativa y ordena mejor la severidad clinica que el amiloide aislado.",
        "El manuscrito debe explicar por que la lectura funcional cambia la interpretacion clinica del declive cognitivo.",
    ],
    img_path=plot_function,
    footer="En el articulo, esta diapositiva debe transformarse en una seccion de progresion funcional.",
)

add_text_image_slide(
    prs,
    "Biomarcadores y ritmo de progresion",
    [
        f"Correlacion centiloid-declive MMSE: r={cor_amyloid_declive:.2f}.",
        f"Correlacion FDG-CDR actual: r={cor_fdg_cdr:.2f}; correlacion hipocampo-CDR actual: r={cor_hippo_cdr:.2f}.",
        "El amiloide senala biologa de enfermedad, mientras FDG e hipocampo se acercan mas a la expresion clinica actual.",
        "La redaccion del articulo debe enlazar biomarcadores de liquido, RM y PET en una narrativa unica.",
    ],
    img_path=plot_biomarkers,
    footer="Esta diapositiva define la seccion multimodal del manuscrito.",
)

add_text_image_slide(
    prs,
    "Tratamiento y sesgo de indicacion",
    [
        f"En enfermedad temprana amiloide positiva, tratados con antiamiloide: {len(treated_early_ad)}; no tratados: {len(untreated_early_ad)}.",
        f"Declive medio MMSE en tratados tempranos: {treated_early_ad['Declive_MMSE_Anual'].mean():.2f}; no tratados: {untreated_early_ad['Declive_MMSE_Anual'].mean():.2f}.",
        "La comparacion cruda por farmaco no estima eficacia: refleja, sobre todo, seleccion por estadio, seguridad y perfil biologico.",
        "El articulo debe discutir explicitamente el sesgo de indicacion antes de interpretar cualquier diferencia entre terapias.",
    ],
    img_path=plot_treatment,
    footer="La seccion del manuscrito debe pasar de comparar tratamientos a interpretar comparadores no aleatorizados.",
)

add_text_image_slide(
    prs,
    "Elegibilidad antiamiloide",
    [
        f"Pacientes candidatos: {antiamyloid_candidates}; proporcion tratados/candidatos: {antiamyloid_treated / antiamyloid_candidates * 100:.1f}%.",
        "La elegibilidad exige amiloide positivo, estadio temprano, ausencia de anticoagulacion, APOE4 no homocigoto y baja carga de microhemorragias.",
        "Los antiamiloides quedan restringidos a fases biologicamente coherentes y radiologicamente seguras.",
        "El manuscrito debe convertir esta diapositiva en una subseccion de seleccion terapeutica basada en biomarcadores y seguridad.",
    ],
    img_path=plot_candidates,
    footer="Esta diapositiva debe expandirse como seccion de elegibilidad antiamiloide.",
)

add_text_image_slide(
    prs,
    "Interpretacion clinica y limitaciones",
    [
        "La cohorte esta construida para docencia: prioriza coherencia clinico-biologica sobre ruido extremo de vida real.",
        "Por eso sirve para entrenar razonamiento por estadio, ATN, seguridad MRI y sesgo de indicacion.",
        "No debe usarse para inferencias regulatorias ni para reclamar magnitudes reales de eficacia.",
        "En el articulo, esta diapositiva debe convertirse en una discusion mas extensa, con limitaciones explicitas.",
    ],
    footer="El manuscrito debe desarrollar esta diapositiva como discusion final y limitaciones.",
)

add_text_image_slide(
    prs,
    "Conclusiones",
    [
        "La lectura conjunta de estadio, cognicion, funcionalidad y biomarcadores produce una cohorte clinicamente interpretable.",
        "La progresion no puede resumirse con una sola variable; exige integrar MMSE, CDR-SB, FAQ y neuroimagen.",
        "La presentacion define el esqueleto del manuscrito: cada apartado debe reescribirse con mas detalle, transiciones y contexto clinico.",
    ],
    footer="Regla docente: cada diapositiva de contenido debe convertirse en una seccion o subseccion del articulo.",
)

output = os.path.join(ROOT, "Ejercicio2_Progresion_Alzheimer.pptx")
prs.save(output)

print(f"Progresores rapidos: {fast}/{total} ({fast / total * 100:.1f}%)")
print("Resumen por estadio:")
print(stage_summary.to_string())
print("\nResumen por tratamiento:")
print(treatment_summary.to_string())
print(f"\nPresentacion generada: {output}")
