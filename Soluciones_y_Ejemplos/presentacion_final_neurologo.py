import pandas as pd
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from neuro_utils import load_and_clean_data

# --- CONFIGURACIÓN ESTÉTICA ---
COLOR_FONDO = RGBColor(229, 232, 232)
COLOR_TITULO = RGBColor(15, 32, 60)
COLOR_TEXTO = RGBColor(40, 40, 40)
PALETA_CLINICA = ["#0F203C", "#2ECC71", "#E74C3C", "#F1C40F"]

sns.set_theme(style="white")
plt.rcParams['axes.facecolor'] = 'none'
plt.rcParams['figure.facecolor'] = 'none'

def add_slide_premium(prs, title_text, content_text="", img_func=None, img_name="temp.png"):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = COLOR_TITULO
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    if content_text:
        body = slide.shapes.placeholders[1]
        body.text = content_text
        for p in body.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = 'Segoe UI'
                r.font.size = Pt(18)

    if img_func:
        img_func(img_name)
        slide.shapes.add_picture(img_name, Inches(5.8), Inches(1.5), width=Inches(3.8))
    return slide

# --- NUEVA FUNCIÓN: PLOT ATROFIA ---
def plot_atrofia(name):
    df = load_and_clean_data()
    plt.figure(figsize=(6, 4))
    sns.regplot(x='MRI_Atrofia', y='Memoria_Episodica', data=df, color="#E74C3C", scatter_kws={'alpha':0.5})
    plt.title("Atrofia vs Rendimiento Cognitivo", color="#0F203C", fontweight='bold')
    plt.tight_layout()
    plt.savefig(name, transparent=True)
    plt.close()

def plot_tratamientos(name):
    df = load_and_clean_data()
    plt.figure(figsize=(6, 4))
    sns.boxplot(x='Tratamiento_Desc', y='EDSS_Actual', data=df, palette=PALETA_CLINICA)
    plt.tight_layout()
    plt.savefig(name, transparent=True)
    plt.close()

# --- EJECUCIÓN ---
def ejecutar_presentacion_final():
    print("Generando Auditoría Clínica EMRR Final...")
    prs = Presentation()
    
    # Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Auditoría Clínica EMRR: Informe Final"
    slide.shapes.placeholders[1].text = "Monitorización de la Progresión y Daño Estructural\nUnidad de NeuroIA 2026"

    # Tratamientos
    add_slide_premium(prs, "1. Eficacia de los Perfiles Terapéuticos", 
                      "- Diferencias significativas en el control de discapacidad.\n- Los inmunomoduladores de alta eficacia muestran una ventaja clara en el EDSS.",
                      img_func=plot_tratamientos, img_name="final_trat.png")

    # ATROFIA (NUEVO)
    add_slide_premium(prs, "2. Atrofia Cortical: El Asesino Silencioso", 
                      "- La pérdida de volumen cortical correlaciona con el deterioro cognitivo.\n- Predictor más fiable de discapacidad irreversible que los brotes puntuales.\n- Requiere monitorización volumétrica anual.",
                      img_func=plot_atrofia, img_name="final_atrofia.png")

    # Conclusiones
    add_slide_premium(prs, "3. Conclusiones Estratégicas", 
                      "1. Prevenir brotes para evitar discapacidad física.\n2. Frenar la atrofia para preservar la función cognitiva.\n3. La IA es la herramienta clave para detectar estos patrones de forma precoz.")

    prs.save("Auditoria_Clinica_EMRR_Completa.pptx")
    print("--- PRESENTACIÓN FINAL GENERADA CON ÉXITO ---")

if __name__ == "__main__":
    ejecutar_presentacion_final()
