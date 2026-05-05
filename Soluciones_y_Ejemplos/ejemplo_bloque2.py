from pptx import Presentation
from neuro_utils import load_and_clean_data

def crear_presentacion_basica():
    df = load_and_clean_data()
    prs = Presentation()
    
    # Diapositiva de Título
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Informe Automático de Cohorte EMRR"
    slide.shapes.placeholders[1].text = f"Análisis de {len(df)} pacientes"
    
    # Diapositiva de Datos
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Resumen de Resultados"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = f"- EDSS Medio: {df['EDSS_Actual'].mean():.2f}"
    p = tf.add_paragraph()
    p.text = f"- Prevalencia OCB+: {df['OCB'].mean()*100:.1f}%"
    
    prs.save("ejemplo_presentacion.pptx")
    print("Presentación 'ejemplo_presentacion.pptx' generada.")

if __name__ == "__main__":
    crear_presentacion_basica()
